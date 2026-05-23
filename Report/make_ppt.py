
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── palette ──────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x0E, 0x27)   # slide bg
CYAN   = RGBColor(0x00, 0xD4, 0xFF)   # accent / headings
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xCC, 0xCC, 0xCC)
GREEN  = RGBColor(0x2E, 0xCC, 0x71)
RED    = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
YELLOW = RGBColor(0xF1, 0xC4, 0x0F)
CARD   = RGBColor(0x13, 0x1A, 0x3A)   # card bg
BORDER = RGBColor(0x00, 0xD4, 0xFF)

blank = prs.slide_layouts[6]  # completely blank

def add_slide():
    return prs.slides.add_slide(blank)

def bg(slide, color=NAVY):
    shp = slide.shapes.add_shape(1, 0, 0, W, H)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def box(slide, l, t, w, h, fill=CARD, alpha=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = BORDER; shp.line.width = Pt(0.8)
    return shp

def txt(slide, text, l, t, w, h, size=18, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

def heading(slide, title):
    # cyan accent bar
    slide.shapes.add_shape(1, Inches(0.4), Inches(0.18),
                           Inches(0.08), Inches(0.65)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = CYAN
    slide.shapes[-1].line.fill.background()
    txt(slide, title, Inches(0.6), Inches(0.12), Inches(12), Inches(0.7),
        size=28, bold=True, color=CYAN)
    # thin horizontal rule
    line = slide.shapes.add_shape(1, Inches(0.4), Inches(0.85),
                                  Inches(12.5), Inches(0.02))
    line.fill.solid(); line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

def img(slide, path, l, t, w, h=None):
    if h:
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        slide.shapes.add_picture(path, l, t, w)

LOGO = r"d:\miniProject\Spam_reporting_portal\Report\iiit manipur.png"
PI   = r"d:\miniProject\Spam_reporting_portal\Report\ppt_imgs"

# ═══════════════════════════════════════════════════════════
# SLIDE 0 – TITLE  (white, matching IIIT format)
# ═══════════════════════════════════════════════════════════
s0 = add_slide()
bg(s0, RGBColor(0xFF,0xFF,0xFF))

txt(s0, "CyberShield: Uncertainty-Aware Fraud Detection\nwith Human-in-the-Loop Calibration for\nCybercrime Intelligence",
    Inches(1), Inches(0.4), Inches(11.3), Inches(2.0),
    size=28, bold=True, color=RGBColor(0x1A,0x1A,0x1A), align=PP_ALIGN.CENTER)

txt(s0, "Presented by", Inches(2.0), Inches(2.55), Inches(4), Inches(0.4),
    size=16, color=RGBColor(0x44,0x44,0x44), align=PP_ALIGN.LEFT)
txt(s0, "ABHISHEK YADAV", Inches(2.0), Inches(2.95), Inches(4), Inches(0.4),
    size=18, bold=True, color=RGBColor(0x0A,0x0E,0x27), align=PP_ALIGN.LEFT)
txt(s0, "Roll No. 230103041\nSemester VI  |  Project – I (CS3201)",
    Inches(2.0), Inches(3.35), Inches(4), Inches(0.7),
    size=15, color=RGBColor(0x33,0x33,0x33), align=PP_ALIGN.LEFT)

txt(s0, "Supervised by", Inches(7.5), Inches(2.55), Inches(4.5), Inches(0.4),
    size=16, color=RGBColor(0x44,0x44,0x44), align=PP_ALIGN.LEFT)
txt(s0, "Dr. Salam Michael Singh", Inches(7.5), Inches(2.95), Inches(4.5), Inches(0.4),
    size=18, bold=True, color=RGBColor(0x0A,0x0E,0x27), align=PP_ALIGN.LEFT)

if os.path.exists(LOGO):
    img(s0, LOGO, Inches(5.4), Inches(3.9), Inches(2.5))

txt(s0, "Department of Computer Science and Engineering\nIndian Institute of Information Technology Senapati, Manipur\n29/04/2026",
    Inches(1), Inches(6.6), Inches(11.3), Inches(0.9),
    size=14, color=RGBColor(0x33,0x33,0x33), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 1 – INTRODUCTION, OBJECTIVE & CONTRIBUTION
# ═══════════════════════════════════════════════════════════
s1 = add_slide()
bg(s1)
heading(s1, "Introduction, Objectives & Unique Contributions")

# Problem card
box(s1, Inches(0.3), Inches(1.0), Inches(4.0), Inches(5.9))
txt(s1, "⚡ THE PROBLEM", Inches(0.4), Inches(1.05), Inches(3.8), Inches(0.45),
    size=13, bold=True, color=CYAN)
problems = [
    "📱 India saw 1.12M+ cybercrime\n   complaints in 2023 (NCRP data)",
    "🔴 Binary spam filters force\n   hard boundaries — no uncertainty",
    "🤖 ML models drift as attackers\n   evolve language & tactics",
    "📋 Victims have no structured path\n   to file a police complaint",
    "🔎 No visibility into repeat\n   malicious URLs / phone numbers",
]
for i, p in enumerate(problems):
    txt(s1, p, Inches(0.4), Inches(1.55)+i*Inches(0.98), Inches(3.7), Inches(0.9),
        size=12, color=LGRAY)

# Contributions card
box(s1, Inches(4.55), Inches(1.0), Inches(8.45), Inches(5.9))
txt(s1, "🏆 UNIQUE CONTRIBUTIONS OF CYBERSHIELD", Inches(4.65), Inches(1.05),
    Inches(8.2), Inches(0.45), size=13, bold=True, color=CYAN)

contribs = [
    (GREEN,  "①", "4-Tier Probabilistic Thresholding",
              "FRAUD / SUSPICIOUS / UNCERTAIN / LEGIT strata\n   replace binary outputs — first in class for SMS fraud"),
    (CYAN,   "②", "Human-in-the-Loop Retraining Loop",
              "UNCERTAIN cases → admin queue → Hard Negative\n   Mining → model re-calibration (concept drift resistance)"),
    (ORANGE, "③", "Automated NCRP PDF Generation",
              "Auto-fills Indian cybercrime complaint with IoCs,\n   probability score & complaint ID — ready to submit"),
    (YELLOW, "④", "Velocity Intelligence Tracker",
              "Real-time freq. DB for malicious URLs & phones;\n   MED/HIGH/CRITICAL threat tiers across submissions"),
    (RED,    "⑤", "Token-Level Word Risk Heatmap",
              "LR coefficient weights highlight fraud-driving\n   tokens per message — explainable AI for citizens"),
]
for i,(col,num,title,desc) in enumerate(contribs):
    y = Inches(1.55) + i*Inches(1.07)
    dot = s1.shapes.add_shape(1, Inches(4.65), y+Inches(0.1),
                               Inches(0.28), Inches(0.28))
    dot.fill.solid(); dot.fill.fore_color.rgb = col
    dot.line.fill.background()
    txt(s1, num, Inches(4.65), y+Inches(0.04), Inches(0.28), Inches(0.3),
        size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(s1, title, Inches(5.0), y, Inches(7.8), Inches(0.35),
        size=13, bold=True, color=col)
    txt(s1, desc,  Inches(5.0), y+Inches(0.34), Inches(7.8), Inches(0.65),
        size=11, color=LGRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 – BACKGROUND / EXISTING WORK
# ═══════════════════════════════════════════════════════════
s2 = add_slide()
bg(s2)
heading(s2, "Background & Existing Work")

# Comparison table
headers = ["System / Paper", "Approach", "Limitation", "CyberShield Fix"]
rows = [
    ["SpamAssassin\n(Apache, 2002)",    "Rule-based filters\n+ Bayesian naive",    "Static rules; breaks\non novel SMS slang",     "Ensemble + HNM\nretraining loop"],
    ["SMS Spam Collection\nBaseline (UCI)",  "Single Naïve Bayes\nTF-IDF",          "Binary only;\nno uncertainty band",      "4-tier probabilistic\nrouting"],
    ["MobiFish\n(Purkait, 2014)",       "URL heuristics\nfor smishing",             "URL-only; ignores\nmessage semantics",       "NLP + URL artifact\ntokenization"],
    ["Deep Spam\n(Roy et al., 2020)",   "LSTM on word\nembeddings",                 "Opaque; no citizen-\nfacing actionability",   "Word Risk Heatmap\n+ NCRP PDF output"],
    ["NCRP Portal\n(Govt. of India)",   "Manual complaint\nentry by victim",        "High friction;\nno automated IoC",          "Pre-filled PDF +\nVelocity Tracker"],
]
col_w = [Inches(2.5), Inches(2.4), Inches(2.8), Inches(2.8)]
col_x = [Inches(0.3), Inches(2.85), Inches(5.3), Inches(8.15)]

# header row
for j,(h,cx,cw) in enumerate(zip(headers,col_x,col_w)):
    b = box(s2, cx, Inches(1.02), cw-Inches(0.06), Inches(0.5), fill=CYAN)
    b.line.fill.background()
    txt(s2, h, cx+Inches(0.05), Inches(1.04), cw-Inches(0.1), Inches(0.45),
        size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

row_colors = [CARD, RGBColor(0x10,0x18,0x35), CARD,
              RGBColor(0x10,0x18,0x35), CARD]
for i, row in enumerate(rows):
    y = Inches(1.55) + i*Inches(1.0)
    for j,(cell,cx,cw) in enumerate(zip(row,col_x,col_w)):
        b = box(s2, cx, y, cw-Inches(0.06), Inches(0.95), fill=row_colors[i])
        col_txt = GREEN if j==3 else (RED if j==2 else WHITE)
        txt(s2, cell, cx+Inches(0.08), y+Inches(0.06), cw-Inches(0.16), Inches(0.85),
            size=11, color=col_txt)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 – PROPOSED SYSTEM / ARCHITECTURE
# ═══════════════════════════════════════════════════════════
s3 = add_slide()
bg(s3)
heading(s3, "Proposed System — CyberShield Architecture")

# Left: pipeline image
pipe_img = os.path.join(PI, "pipeline_flow.png")
if os.path.exists(pipe_img):
    img(s3, pipe_img, Inches(0.3), Inches(1.0), Inches(6.8), Inches(5.7))

# Right: details
box(s3, Inches(7.3), Inches(1.0), Inches(5.7), Inches(5.9))
txt(s3, "🛠 TOOLS & STACK", Inches(7.45), Inches(1.08), Inches(5.4), Inches(0.4),
    size=13, bold=True, color=CYAN)

stack = [
    ("Backend / API",  "Python 3.11 · Flask · SQLAlchemy"),
    ("ML Engine",      "scikit-learn: LR + SVM + Naïve Bayes\n25 000-feature FeatureUnion TF-IDF"),
    ("Auth",           "Google OAuth 2.0 (admin gating)"),
    ("Database",       "SQLite: complaints · velocity · HITL queue"),
    ("PDF Output",     "ReportLab — NCRP-compliant"),
    ("Explainability", "LR coeff. weights → Word Risk Heatmap"),
]
for i,(label,val) in enumerate(stack):
    y = Inches(1.55) + i*Inches(0.77)
    txt(s3, label, Inches(7.45), y, Inches(2.0), Inches(0.35),
        size=11, bold=True, color=CYAN)
    txt(s3, val,   Inches(9.5),  y, Inches(3.35), Inches(0.65),
        size=11, color=WHITE)

txt(s3, "⚙ PIPELINE — 10 STAGES", Inches(7.45), Inches(6.1), Inches(5.4), Inches(0.4),
    size=11, bold=True, color=CYAN)
txt(s3, "Input → Preprocess → TF-IDF → Ensemble Score →\nThreshold Router → DB Log → Velocity Update →\nHITL Queue (if UNCERTAIN) → NCRP PDF → Response",
    Inches(7.45), Inches(6.5), Inches(5.4), Inches(0.85),
    size=10, color=LGRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 – RESULTS & ANALYSIS
# ═══════════════════════════════════════════════════════════
s4 = add_slide()
bg(s4)
heading(s4, "Results & Analysis")

# ── KPI cards row ──
kpis = [
    ("94.85%",  "Test Accuracy",    GREEN),
    ("0.9899",  "ROC-AUC",          CYAN),
    ("26,534",  "Labeled Samples",  ORANGE),
    ("< 500ms", "Inference Latency",YELLOW),
    ("3,980",   "Frozen Test Set",  RGBColor(0xAF,0x7A,0xC5)),
]
kw = Inches(2.4)
for i,(val,lab,col) in enumerate(kpis):
    x = Inches(0.25) + i*(kw+Inches(0.08))
    b = box(s4, x, Inches(1.0), kw, Inches(1.15), fill=CARD)
    txt(s4, val, x, Inches(1.08), kw, Inches(0.55),
        size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s4, lab, x, Inches(1.62), kw, Inches(0.42),
        size=11, color=LGRAY, align=PP_ALIGN.CENTER)

# ── HNM ablation mini-table ──
box(s4, Inches(0.25), Inches(2.35), Inches(4.6), Inches(4.5))
txt(s4, "📊 HNM Ablation Table", Inches(0.4), Inches(2.42),
    Inches(4.3), Inches(0.38), size=12, bold=True, color=CYAN)
abl_h = ["Metric","Pre-HNM","Post-HNM","Δ"]
abl_r = [
    ["Accuracy",       "0.9482","0.9485","+0.0003"],
    ["Precision",      "0.9508","0.9494","−0.0014"],
    ["Recall (Fraud)", "0.9419","0.9440","+0.0021"],
    ["Macro F1",       "0.9482","0.9484","+0.0002"],
    ["ROC-AUC",        "0.9893","0.9899","+0.0006"],
]
cws2 = [Inches(1.6),Inches(0.9),Inches(1.0),Inches(0.85)]
cxs2 = [Inches(0.3),Inches(1.95),Inches(2.88),Inches(3.92)]
for j,(h,cx,cw) in enumerate(zip(abl_h,cxs2,cws2)):
    b2 = box(s4, cx, Inches(2.82), cw-Inches(0.04), Inches(0.38), fill=CYAN)
    b2.line.fill.background()
    txt(s4, h, cx+Inches(0.04), Inches(2.85), cw, Inches(0.32),
        size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
for i,row in enumerate(abl_r):
    ry = Inches(3.23)+i*Inches(0.62)
    for j,(cell,cx,cw) in enumerate(zip(row,cxs2,cws2)):
        rc = CARD if i%2==0 else RGBColor(0x10,0x18,0x35)
        box(s4, cx, ry, cw-Inches(0.04), Inches(0.58), fill=rc)
        col_c = GREEN if (j==3 and cell.startswith("+")) else (RED if (j==3 and cell.startswith("−")) else WHITE)
        txt(s4, cell, cx+Inches(0.04), ry+Inches(0.06), cw, Inches(0.45),
            size=10, color=col_c, align=PP_ALIGN.CENTER)

# ── ROC/PR curves image ──
roc_img = os.path.join(PI,"fig_roc_pr_curves.png")
if os.path.exists(roc_img):
    img(s4, roc_img, Inches(5.05), Inches(2.35), Inches(4.05), Inches(2.15))

# ── Threshold sensitivity ──
thr_img = os.path.join(PI,"fig_threshold_sensitivity.png")
if os.path.exists(thr_img):
    img(s4, thr_img, Inches(9.15), Inches(2.35), Inches(3.9), Inches(2.15))

# ── Confusion matrix ──
cm_img = os.path.join(PI,"fig_confusion_matrices.png")
if os.path.exists(cm_img):
    img(s4, cm_img, Inches(5.05), Inches(4.6), Inches(8.0), Inches(2.3))

# ── Threshold caption ──
txt(s4, "Threshold Sensitivity with FRAUD/SUSPICIOUS/UNCERTAIN/LEGIT risk strata (Fig 6.3)",
    Inches(9.15), Inches(4.55), Inches(3.9), Inches(0.5), size=9, color=LGRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 – CONCLUSION
# ═══════════════════════════════════════════════════════════
s5 = add_slide()
bg(s5)
heading(s5, "Conclusion & Future Scope")

# Summary cards
summ = [
    (GREEN,  "✅ Achieved",
     "94.85% accuracy · ROC-AUC 0.9899 on 3,980 frozen test samples across\n26,534 labeled records from 8 diverse datasets"),
    (CYAN,   "🎯 Unique Value",
     "4-tier uncertainty routing + HITL loop is absent in all surveyed systems.\nHNM improved fraud recall by +0.21pp with zero manual labeling cost"),
    (ORANGE, "📋 Actionable Output",
     "First system to auto-generate NCRP-compliant PDF complaints + Velocity\nIntelligence DB for law-enforcement threat tracking"),
    (YELLOW, "⚡ Efficiency",
     "End-to-end inference < 500ms on commodity hardware (Intel i5, 8 GB RAM).\nAll 5 functional + 5 non-functional requirements verified"),
]
for i,(col,label,desc) in enumerate(summ):
    y = Inches(1.05) + i*Inches(1.2)
    dot = s5.shapes.add_shape(1, Inches(0.3), y+Inches(0.12),
                               Inches(0.22), Inches(0.22))
    dot.fill.solid(); dot.fill.fore_color.rgb = col
    dot.line.fill.background()
    txt(s5, label, Inches(0.6), y, Inches(2.3), Inches(0.4),
        size=13, bold=True, color=col)
    txt(s5, desc, Inches(0.6), y+Inches(0.4), Inches(7.5), Inches(0.72),
        size=11, color=LGRAY)

# Future scope box
box(s5, Inches(8.5), Inches(1.0), Inches(4.55), Inches(5.5))
txt(s5, "🚀 FUTURE SCOPE", Inches(8.65), Inches(1.08), Inches(4.3), Inches(0.4),
    size=13, bold=True, color=CYAN)
future = [
    ("Near-Term",  "Direct NCRP REST API integration\n— eliminate manual PDF upload"),
    ("Near-Term",  "OTP/QR Code artifact extraction\n— expand IoC coverage"),
    ("Mid-Term",   "Multilingual NLP pipeline\n— Devanagari, Manipuri scripts"),
    ("Mid-Term",   "Transformer embeddings (BERT)\n— semantic fraud detection"),
    ("Long-Term",  "Federated learning across state\n— cyber cells (privacy-safe)"),
    ("Long-Term",  "Real-time threat intelligence\nfeed API for law enforcement"),
]
tags = {
    "Near-Term": RGBColor(0x2E,0xCC,0x71),
    "Mid-Term":  ORANGE,
    "Long-Term": RED,
}
for i,(tag,desc) in enumerate(future):
    y = Inches(1.55) + i*Inches(0.82)
    tb = s5.shapes.add_shape(1, Inches(8.65), y, Inches(1.05), Inches(0.25))
    tb.fill.solid(); tb.fill.fore_color.rgb = tags[tag]
    tb.line.fill.background()
    txt(s5, tag, Inches(8.65), y-Inches(0.01), Inches(1.1), Inches(0.28),
        size=8, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(s5, desc, Inches(9.75), y, Inches(3.15), Inches(0.7),
        size=11, color=WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 – REFERENCES
# ═══════════════════════════════════════════════════════════
s6 = add_slide()
bg(s6)
heading(s6, "References")

refs = [
    "[1] Almeida, T.A., Hidalgo, J.M.G., Yamakami, A. (2011). Contributions to the study of SMS spam filtering. ACM DocEng, pp. 259–262.",
    "[2] Settles, B. (2012). Active Learning. Synthesis Lectures on AI & ML, Morgan & Claypool.",
    "[3] Holzinger, A. (2016). Interactive machine learning for health informatics. Brain Informatics, 3(2), 119–131.",
    "[4] Chawla, N.V. et al. (2002). SMOTE: Synthetic Minority Over-sampling Technique. JAIR, 16, 321–357.",
    "[5] Roy, S., et al. (2020). A Deep Learning based Artificial Intelligence Model for SMS Spam Detection. IJEAT, 9(3).",
    "[6] Pedregosa, F. et al. (2011). Scikit-learn: Machine learning in Python. JMLR, 12, 2825–2830.",
    "[7] Purkait, S. (2014). Phishing counter measures and their effectiveness. Information Management & Computer Security, 22(5).",
    "[8] Ministry of Home Affairs, Govt. of India (2023). Annual Report on Cybercrime — NCRP Statistics 2023.",
    "[9] Friedl, J.E.F. (2006). Mastering Regular Expressions, 3rd ed. O'Reilly Media.",
    "[10] National Crime Records Bureau (2023). Crime in India Report — Cyber Crimes Chapter.",
]
for i, ref in enumerate(refs):
    y = Inches(1.02) + i*Inches(0.55)
    txt(s6, ref, Inches(0.4), y, Inches(12.5), Inches(0.52), size=11, color=LGRAY)

out = r"d:\miniProject\Spam_reporting_portal\CyberShield_Presentation.pptx"
prs.save(out)
print("Saved:", out)
