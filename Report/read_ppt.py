
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import re

prs = Presentation(r'd:\miniProject\Spam_reporting_portal\Report\Smart UPI Fraud Detection.pptx')
print(f'Slides: {len(prs.slides)}  W={prs.slide_width}  H={prs.slide_height}')

def try_rgb(color_obj):
    try: return str(color_obj.rgb)
    except: return None

for si, slide in enumerate(prs.slides):
    print(f'\n=== SLIDE {si} ===')
    for shape in slide.shapes:
        name = shape.name
        # solid fills
        try:
            f = shape.fill
            if str(f.type) == 'SOLID (1)':
                rgb = try_rgb(f.fore_color)
                if rgb: print(f'  [SOLID] {name}: {rgb}')
        except: pass

        # gradient fills via XML
        try:
            xml = shape._element.xml
            grad_stops = re.findall(r'<a:gs pos="(\d+)".*?<a:srgbClr val="([0-9A-Fa-f]{6})"', xml, re.DOTALL)
            for pos, col in grad_stops:
                print(f'  [GRAD]  {name}: pos={int(pos)/1000:.1f}% color=#{col}')
        except: pass

        # text
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text.strip()
                    if not t: continue
                    rgb = try_rgb(run.font.color)
                    sz  = run.font.size
                    bd  = run.font.bold
                    sz_pt = round(sz.pt, 1) if sz else None
                    print(f'  [TEXT]  "{t[:50]}" sz={sz_pt} bold={bd} color={rgb}')

# Also dump theme colors
print('\n=== THEME COLORS ===')
try:
    theme = prs.slide_master.element.find('.//' + qn('a:theme'))
    if theme is None:
        theme_xml = prs.slide_master.element.xml
        colors = re.findall(r'<a:srgbClr val="([0-9A-Fa-f]{6})"', theme_xml)
        for c in sorted(set(colors)):
            print(f'  #{c}')
except Exception as e:
    print(f'Theme error: {e}')
