
# CyberShield PPT – Academic Content + UPI Palette + Rounded Rectangles + Larger Custom Font
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

# ── exact UPI PPT palette ──────────────────────────────────
BG      = RGBColor(0xF9,0xFA,0xFB)
NAVY    = RGBColor(0x1E,0x3A,0x5F)
TEAL    = RGBColor(0x2E,0x8B,0x7B)
BLUE    = RGBColor(0x3B,0x82,0xF6)
GREEN   = RGBColor(0x10,0xB9,0x81)
PURPLE  = RGBColor(0x8B,0x5C,0xF6)
RED     = RGBColor(0xEF,0x44,0x44)
AMBER   = RGBColor(0xF5,0x9E,0x0B)
GRAY    = RGBColor(0x6B,0x72,0x80)
DARK    = RGBColor(0x1E,0x29,0x3B)
WHITE   = RGBColor(0xFF,0xFF,0xFF)

C_BLUE  = RGBColor(0xDB,0xEA,0xFE)
C_GREEN = RGBColor(0xD1,0xFA,0xE5)
C_PURP  = RGBColor(0xED,0xE9,0xFE)
C_AMBER = RGBColor(0xFE,0xF3,0xC7)
C_RED   = RGBColor(0xFF,0xED,0xD5)
C_GRAY  = RGBColor(0xE5,0xE7,0xEB)

# Base Font
FONT_NAME = "Segoe UI"

blank = prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(blank)

def rect(sl, l,t,w,h, fill=BG, line_col=None, line_w=Pt(0), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = sl.shapes.add_shape(shape_type,l,t,w,h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_col:
        s.line.color.rgb = line_col; s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def tb(sl, text, l,t,w,h, sz=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    box = sl.shapes.add_textbox(l,t,w,h)
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = FONT_NAME
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return box

def hdr(sl, title):
    rect(sl, Inches(0.35), Inches(0.18), Inches(0.1), Inches(0.6), fill=TEAL, shape_type=MSO_SHAPE.RECTANGLE)
    tb(sl, title, Inches(0.55), Inches(0.14), Inches(12.2), Inches(0.7), sz=30, bold=True, color=NAVY)
    rect(sl, Inches(0.35), Inches(0.82), Inches(12.6), Pt(2), fill=TEAL, shape_type=MSO_SHAPE.RECTANGLE)

def card(sl, l,t,w,h, fill=C_BLUE, border=BLUE):
    return rect(sl,l,t,w,h, fill=fill, line_col=border, line_w=Pt(1), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

def img(sl, path, l,t,w,h=None):
    if os.path.exists(path):
        if h: sl.shapes.add_picture(path,l,t,w,h)
        else: sl.shapes.add_picture(path,l,t,w)

LOGO = r"d:\miniProject\Spam_reporting_portal\Report\iiit manipur.png"
PI   = r"d:\miniProject\Spam_reporting_portal\Report\ppt_imgs"

# ════════════════════════════════════════════════════════════
# TITLE SLIDE 
# ════════════════════════════════════════════════════════════
s0 = slide()
rect(s0,0,0,W,H, fill=WHITE, shape_type=MSO_SHAPE.RECTANGLE)
rect(s0, 0, 0, W, Inches(0.12), fill=TEAL, shape_type=MSO_SHAPE.RECTANGLE)
rect(s0, 0, H-Inches(0.12), W, Inches(0.12), fill=TEAL, shape_type=MSO_SHAPE.RECTANGLE)
rect(s0, Inches(2.5), Inches(2.22), Inches(8.3), Pt(2.5), fill=TEAL, shape_type=MSO_SHAPE.RECTANGLE)
rect(s0, Inches(2.5), Inches(2.32), Inches(8.3), Pt(1),   fill=C_GREEN, shape_type=MSO_SHAPE.RECTANGLE)

tb(s0, "CyberShield: Uncertainty-Aware Fraud Detection", Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.85), sz=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
tb(s0, "with Human-in-the-Loop Calibration for Cybercrime Intelligence", Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.6), sz=24, color=TEAL, align=PP_ALIGN.CENTER)
tb(s0, "A report submitted for the course Project – I (CS3201)", Inches(0.5), Inches(1.82), Inches(12.3), Inches(0.4), sz=16, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

tb(s0, "Submitted By",   Inches(1.8), Inches(2.55), Inches(4), Inches(0.38), sz=18, bold=True, color=DARK)
tb(s0, "ABHISHEK YADAV", Inches(1.8), Inches(2.95), Inches(4), Inches(0.4), sz=22, bold=True, color=NAVY)
tb(s0, "Roll No. 230103041  |  Semester VI", Inches(1.8), Inches(3.38), Inches(4), Inches(0.38), sz=16, color=GRAY)

rect(s0, Inches(6.55), Inches(2.55), Pt(1.5), Inches(1.3), fill=C_GRAY, shape_type=MSO_SHAPE.RECTANGLE)

tb(s0, "Supervised By", Inches(6.9), Inches(2.55), Inches(5), Inches(0.38), sz=18, bold=True, color=DARK)
tb(s0, "Dr. Salam Michael Singh", Inches(6.9), Inches(2.95), Inches(5), Inches(0.4), sz=22, bold=True, color=NAVY)

if os.path.exists(LOGO): img(s0, LOGO, Inches(5.55), Inches(3.95), Inches(2.2))

tb(s0,"Department of Computer Science and Engineering", Inches(0.5),Inches(6.35),Inches(12.3),Inches(0.35), sz=15,color=GRAY,align=PP_ALIGN.CENTER)
tb(s0,"Indian Institute of Information Technology Senapati, Manipur", Inches(0.5),Inches(6.7),Inches(12.3),Inches(0.35), sz=15,color=GRAY,align=PP_ALIGN.CENTER)
tb(s0,"April, 2026", Inches(0.5),Inches(7.05),Inches(12.3),Inches(0.35), sz=16,bold=True,color=TEAL,align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 1 – INTRODUCTION, OBJECTIVE & CONTRIBUTION
# ════════════════════════════════════════════════════════════
s1 = slide()
rect(s1,0,0,W,H,fill=BG, shape_type=MSO_SHAPE.RECTANGLE)
hdr(s1,"Introduction, Objective & Contribution")

card(s1, Inches(0.3), Inches(1.0), Inches(5.8), Inches(6.2), C_BLUE, BLUE)
tb(s1,"PROBLEM DEFINITION", Inches(0.45),Inches(1.08),Inches(5.5),Inches(0.42), sz=15,bold=True,color=BLUE)
probs=[
    ("Limitation of Binary Classification", "Standard ML forces binary buckets (Fraud vs. Legit), failing on ambiguous phrasing, leading to high false-positive rates."),
    ("Semantic Evading & Concept Drift", "Attackers leverage urgency and financial requests, constantly shifting vocabulary to evade static signature detection."),
    ("Fragmented Reporting Workflows", "Victims lack automated tools to extract Indicators of Compromise (IoCs) and format actionable legal complaints for the NCRP."),
]
for i,(t,p) in enumerate(probs):
    rect(s1,Inches(0.45),Inches(1.68)+i*Inches(1.45), Inches(0.12),Inches(0.12),fill=BLUE, shape_type=MSO_SHAPE.OVAL)
    tb(s1,t,Inches(0.65),Inches(1.63)+i*Inches(1.45), Inches(5.2),Inches(0.3), sz=14,bold=True,color=NAVY)
    tb(s1,p,Inches(0.65),Inches(1.88)+i*Inches(1.45), Inches(5.2),Inches(0.8), sz=13,color=DARK)

card(s1,Inches(6.3),Inches(1.0),Inches(6.7),Inches(6.2),C_GREEN,GREEN)
tb(s1,"CYBERSHIELD: PRIMARY CONTRIBUTIONS",Inches(6.45),Inches(1.08),Inches(6.4),Inches(0.42), sz=15,bold=True,color=GREEN)
contribs=[
    (BLUE,  "Uncertainty-Aware Classification", "Probability-bounded risk strata (FRAUD, SUSPICIOUS, UNCERTAIN, LEGIT) replacing forced binary predictions."),
    (GREEN, "Human-in-the-Loop Calibration", "Administrative routing for uncertain data (0.30 \u2264 P < 0.70) coupled with Hard Negative Mining for dynamic recalibration."),
    (PURPLE,"Token-Level Explainability", "Word Risk Heatmaps utilizing model coefficient weights to visually highlight specific semantic tokens driving predictions."),
    (AMBER, "Automated NCRP Adjudication", "One-click generation of NCRP-compliant PDFs containing extracted IoCs, ready for portal submission."),
    (RED,   "Velocity Intelligence Tracking", "Database tracking the frequency velocity of malicious URLs and phone numbers to identify organized threat vectors."),
]
for i,(col,title,desc) in enumerate(contribs):
    y=Inches(1.58)+i*Inches(0.95)
    dot=s1.shapes.add_shape(MSO_SHAPE.OVAL,Inches(6.45),y+Inches(0.08), Inches(0.15),Inches(0.15))
    dot.fill.solid();dot.fill.fore_color.rgb=col;dot.line.fill.background()
    tb(s1,title,Inches(6.7),y,Inches(6.1),Inches(0.3), sz=14,bold=True,color=NAVY)
    tb(s1,desc, Inches(6.7),y+Inches(0.25),Inches(6.1),Inches(0.6), sz=13,color=DARK)

# ════════════════════════════════════════════════════════════
# SLIDE 2 – BACKGROUND / EXISTING WORK
# ════════════════════════════════════════════════════════════
s2=slide()
rect(s2,0,0,W,H,fill=BG, shape_type=MSO_SHAPE.RECTANGLE)
hdr(s2,"Background / Existing Work")

headers=["Existing Work / System", "Approach Used", "Primary Limitations", "Proposed Solution (CyberShield)"]
rows=[
    ["SpamAssassin", "Rule-based filtering with\nBayesian Naive integration", "Static rules become obsolete rapidly\nagainst novel semantic phrasing", "Ensemble architecture with HNM\nretraining for concept drift resistance"],
    ["SMS Spam Collection Baseline", "Single Naive Bayes on\nTF-IDF feature vectors", "Forces binary outcomes;\nlacks uncertainty acknowledgment", "4-tier probabilistic routing\nbased on confidence thresholds"],
    ["MobiFish (Purkait, 2014)", "URL-specific heuristics\nfor smishing detection", "Focuses solely on URLs, ignoring\npersuasive natural language text", "FeatureUnion of NLP tokens and\nlexical artifact extraction"],
    ["Deep Spam (Roy et al., 2020)", "LSTM on word embeddings\n(Deep Learning)", "Opaque \"black-box\" model; no\ncitizen-facing explainability", "Word Risk Heatmap generation\nusing linear coefficient weights"],
    ["NCRP Portal (Govt. of India)", "Manual, unguided complaint\nentry by the victim", "High friction workflow; users struggle\nto identify technical IoCs", "Automated pre-filled PDF generation\nwith extracted intelligence"],
]
cws=[Inches(2.5),Inches(2.7),Inches(3.3),Inches(3.8)]
cxs=[Inches(0.3),Inches(2.9),Inches(5.7),Inches(9.1)]
h_fills=[NAVY,BLUE,RED,GREEN]
for j,(h,cx,cw,hf) in enumerate(zip(headers,cxs,cws,h_fills)):
    rect(s2,cx,Inches(1.05),cw-Inches(0.1),Inches(0.48),fill=hf, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s2,h,cx,Inches(1.15),cw-Inches(0.1),Inches(0.42), sz=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

row_bg=[C_BLUE,BG,C_GREEN,BG,C_AMBER]
for i,row in enumerate(rows):
    y=Inches(1.57)+i*Inches(0.97)
    for j,(cell,cx,cw) in enumerate(zip(row,cxs,cws)):
        rect(s2,cx,y,cw-Inches(0.1),Inches(0.9),fill=row_bg[i], line_col=C_GRAY,line_w=Pt(0.5), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        tc=GREEN if j==3 else (RED if j==2 else DARK)
        tb(s2,cell,cx+Inches(0.1),y+Inches(0.15),cw-Inches(0.2), Inches(0.78),sz=12,color=tc)

# ════════════════════════════════════════════════════════════
# SLIDE 3 – PROPOSED SYSTEM / ARCHITECTURE
# ════════════════════════════════════════════════════════════
s3=slide()
rect(s3,0,0,W,H,fill=BG, shape_type=MSO_SHAPE.RECTANGLE)
hdr(s3,"Proposed System / Architecture")

pipe=os.path.join(PI,"pipeline_flow.png")
if os.path.exists(pipe): img(s3,pipe,Inches(0.3),Inches(1.0),Inches(7.2),Inches(5.8))

card(s3,Inches(7.7),Inches(1.0),Inches(5.3),Inches(5.8),C_BLUE,BLUE)
tb(s3,"SYSTEM DESIGN & TOOLS",Inches(7.85),Inches(1.08),Inches(5.0),Inches(0.4), sz=15,bold=True,color=BLUE)
stack=[
    (NAVY,  "Text Vectorization", "FeatureUnion combining word n-grams (1,2) and char n-grams (2,5) into a 25,000-dimensional TF-IDF matrix."),
    (BLUE,  "Classification Engine", "Mathematical ensemble of Logistic Regression, Support Vector Machines (Platt Scaled), and Naive Bayes."),
    (GREEN, "Threshold Routing", "P \u2265 0.90 (FRAUD): Auto PDF\n0.30 \u2264 P < 0.70 (UNCERTAIN): HITL Queue\nP < 0.30 (LEGIT): Cleared"),
    (PURPLE,"Tech Stack", "Python 3.11, Flask, Scikit-Learn, SQLite, ReportLab."),
]
for i,(col,label,val) in enumerate(stack):
    y=Inches(1.55)+i*Inches(1.05)
    dot2=s3.shapes.add_shape(MSO_SHAPE.OVAL,Inches(7.85),y+Inches(0.06), Inches(0.15),Inches(0.15))
    dot2.fill.solid();dot2.fill.fore_color.rgb=col;dot2.line.fill.background()
    tb(s3,label,Inches(8.1),y,Inches(4.8),Inches(0.3),sz=14,bold=True,color=col)
    tb(s3,val,  Inches(8.1), y+Inches(0.25),Inches(4.8),Inches(0.7),sz=13,color=DARK)

rect(s3,Inches(7.7),Inches(5.5),Inches(5.3),Pt(1.5),fill=TEAL, shape_type=MSO_SHAPE.RECTANGLE)
tb(s3,"PROCESS WORKFLOW",Inches(7.85),Inches(5.6),Inches(5.0),Inches(0.35), sz=13,bold=True,color=NAVY)
tb(s3,"Input \u2192 NLP Preprocessing \u2192 TF-IDF Transformation \u2192 Ensemble Scoring \u2192 Threshold Evaluation \u2192 Velocity Update \u2192 Output Generation.", Inches(7.85),Inches(5.85),Inches(5.0),Inches(0.7),sz=12,color=DARK)

# ════════════════════════════════════════════════════════════
# SLIDE 4 – RESULT & ANALYSIS
# ════════════════════════════════════════════════════════════
s4=slide()
rect(s4,0,0,W,H,fill=BG, shape_type=MSO_SHAPE.RECTANGLE)
hdr(s4,"Result & Analysis")

kpis=[
    ("94.85%","Test Accuracy",   C_GREEN,GREEN),
    ("0.9899", "ROC-AUC",        C_BLUE, BLUE),
    ("26,534", "Labeled Dataset",C_PURP,PURPLE),
    ("3,980",  "Frozen Test Set", C_RED, RED),
]
kw=Inches(3.05)
for i,(val,lab,bg2,col2) in enumerate(kpis):
    x=Inches(0.3)+i*(kw+Inches(0.15))
    card(s4,x,Inches(1.0),kw,Inches(1.0),bg2,col2)
    tb(s4,val,x,Inches(1.05),kw,Inches(0.5),sz=28,bold=True, color=col2,align=PP_ALIGN.CENTER)
    tb(s4,lab,x,Inches(1.55),kw,Inches(0.4),sz=13,color=DARK, align=PP_ALIGN.CENTER)

card(s4,Inches(0.3),Inches(2.1),Inches(4.6),Inches(4.8),C_GREEN,GREEN)
tb(s4,"Hard Negative Mining (Ablation)",Inches(0.45),Inches(2.15),Inches(4.3),Inches(0.38), sz=14,bold=True,color=GREEN)
abl_h=["Metric","Pre-HNM","Post-HNM","Delta"]
abl_r=[
    ["Accuracy",      "0.9482","0.9485","+0.0003"],
    ["Precision",     "0.9508","0.9494","-0.0014"],
    ["Recall(Fraud)", "0.9419","0.9440","+0.0021"],
    ["Macro F1",      "0.9482","0.9484","+0.0002"],
    ["ROC-AUC",       "0.9893","0.9899","+0.0006"],
]
cws2=[Inches(1.55),Inches(0.88),Inches(1.0),Inches(0.88)]
cxs2=[Inches(0.35),Inches(1.95),Inches(2.87),Inches(3.91)]
for j,(h,cx,cw) in enumerate(zip(abl_h,cxs2,cws2)):
    hf2=[NAVY,BLUE,GREEN,TEAL][j]
    rect(s4,cx,Inches(2.55),cw-Inches(0.04),Inches(0.35),fill=hf2, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s4,h,cx,Inches(2.6),cw,Inches(0.3), sz=11,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
for i,row in enumerate(abl_r):
    ry=Inches(2.95)+i*Inches(0.7)
    rbg=C_BLUE if i%2==0 else WHITE
    for j,(cell,cx,cw) in enumerate(zip(row,cxs2,cws2)):
        rect(s4,cx,ry,cw-Inches(0.04),Inches(0.65),fill=rbg, line_col=C_GRAY,line_w=Pt(0.5), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        tc2=GREEN if(j==3 and cell.startswith("+")) else(RED if(j==3 and cell.startswith("-")) else DARK)
        tb(s4,cell,cx,ry+Inches(0.15),cw,Inches(0.4), sz=12,color=tc2,align=PP_ALIGN.CENTER)

roc=os.path.join(PI,"fig_roc_pr_curves.png")
if os.path.exists(roc): img(s4,roc,Inches(5.1),Inches(2.2),Inches(4.0),Inches(2.25))
thr=os.path.join(PI,"fig_threshold_sensitivity.png")
if os.path.exists(thr): img(s4,thr,Inches(9.2),Inches(2.2),Inches(3.9),Inches(2.25))
cm=os.path.join(PI,"fig_confusion_matrices.png")
if os.path.exists(cm): img(s4,cm,Inches(5.1),Inches(4.6),Inches(8.0),Inches(2.4))

# ════════════════════════════════════════════════════════════
# SLIDE 5 – CONCLUSION
# ════════════════════════════════════════════════════════════
s5=slide()
rect(s5,0,0,W,H,fill=BG, shape_type=MSO_SHAPE.RECTANGLE)
hdr(s5,"Conclusion & Future Scope")

card(s5, Inches(0.3), Inches(1.0), Inches(7.5), Inches(5.9), C_BLUE, BLUE)
tb(s5,"SUMMARY OF WORK DONE", Inches(0.5),Inches(1.1),Inches(7.0),Inches(0.4), sz=15,bold=True,color=BLUE)

summ=[
    ("Mathematical Efficacy", "The calibrated ensemble over a 25,000-feature TF-IDF matrix achieved 94.85% test-set accuracy and 0.9899 ROC-AUC."),
    ("System Safety via Uncertainty", "The wide UNCERTAIN band (0.30 \u2264 P < 0.70) acts as a deliberate safety mechanism against catastrophic false negatives, proving that acknowledging uncertainty enhances, rather than weakens, AI safety in high-stakes environments."),
    ("Bridging the Intelligence Gap", "By intertwining natural language semantic analysis with lexical Velocity Tracking and automated NCRP PDF generation, the project successfully bridges the gap between raw ML predictions and actionable law-enforcement intelligence."),
    ("Operational Efficiency", "The system achieves end-to-end inference latency under 500ms on commodity hardware, satisfying all non-functional constraints.")
]
for i,(label,desc) in enumerate(summ):
    y=Inches(1.6)+i*Inches(1.05)
    dot=s5.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.5),y+Inches(0.06), Inches(0.12),Inches(0.12))
    dot.fill.solid();dot.fill.fore_color.rgb=NAVY;dot.line.fill.background()
    tb(s5,label,Inches(0.7),y,Inches(6.8),Inches(0.3), sz=14,bold=True,color=NAVY)
    tb(s5,desc, Inches(0.7),y+Inches(0.25),Inches(6.8),Inches(0.7), sz=13,color=DARK)

card(s5,Inches(8.0),Inches(1.0),Inches(5.0),Inches(5.9),C_PURP,PURPLE)
tb(s5,"FUTURE SCOPE",Inches(8.2),Inches(1.1),Inches(4.6),Inches(0.4), sz=15,bold=True,color=PURPLE)
future=[
    (GREEN, "Near-Term", "Direct REST API Integration with NCRP to eliminate manual PDF upload requirements."),
    (GREEN, "Near-Term", "Expansion of Artifact Extraction to include OTP/QR Code payload heuristics."),
    (AMBER, "Mid-Term", "Implementation of a Multilingual NLP pipeline (e.g., Devanagari, Meitei Mayek)."),
    (AMBER, "Mid-Term", "Integration of Transformer-based semantic embeddings (e.g., BERT/RoBERTa)."),
    (RED,   "Long-Term", "Deployment of Federated Learning architectures across state cyber-cells to preserve privacy."),
]
for i,(col,tag,desc) in enumerate(future):
    y=Inches(1.6)+i*Inches(0.85)
    rect(s5,Inches(8.2),y,Inches(1.1),Inches(0.28),fill=col, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s5,tag,Inches(8.2),y-Inches(0.02),Inches(1.1),Inches(0.3), sz=10,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    tb(s5,desc,Inches(9.4),y,Inches(3.4),Inches(0.8),sz=13,color=DARK)

# ════════════════════════════════════════════════════════════
# SLIDE 6 – REFERENCES
# ════════════════════════════════════════════════════════════
s6=slide()
rect(s6,0,0,W,H,fill=BG, shape_type=MSO_SHAPE.RECTANGLE)
hdr(s6,"References")
refs=[
    "[1] Almeida, T.A., Hidalgo, J.M.G., & Yamakami, A. (2011). Contributions to the study of SMS spam filtering. ACM DocEng, 259-262.",
    "[2] Settles, B. (2012). Active Learning. Synthesis Lectures on AI and Machine Learning. Morgan & Claypool.",
    "[3] Holzinger, A. (2016). Interactive machine learning for health informatics. Brain Informatics, 3(2), 119-131.",
    "[4] Chawla, N.V. et al. (2002). SMOTE: Synthetic Minority Over-sampling Technique. JAIR, 16, 321-357.",
    "[5] Roy, S. et al. (2020). A Deep Learning based AI Model for SMS Spam Detection. IJEAT, 9(3).",
    "[6] Pedregosa, F. et al. (2011). Scikit-learn: Machine Learning in Python. JMLR, 12, 2825-2830.",
    "[7] Purkait, S. (2014). Phishing counter measures and their effectiveness. Info. Mgmt. & Comp. Security, 22(5).",
    "[8] Ministry of Home Affairs, Govt. of India (2023). Annual Report on Cybercrime - NCRP Statistics 2023.",
]
for i,ref in enumerate(refs):
    card(s6,Inches(0.3),Inches(1.02)+i*Inches(0.7), Inches(12.6),Inches(0.6), C_BLUE if i%2==0 else WHITE, C_GRAY)
    tb(s6,ref,Inches(0.45),Inches(1.15)+i*Inches(0.7), Inches(12.3),Inches(0.5),sz=13,color=DARK)

out = r"d:\miniProject\Spam_reporting_portal\CyberShield_Presentation_Final.pptx"
prs.save(out)
print("Saved:", out)
