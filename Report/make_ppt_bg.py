
# CyberShield PPT – Academic Content + Custom Background + Rounded Rectangles + Larger Custom Font
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

BG_IMG = r"d:\miniProject\Spam_reporting_portal\Report\bg.png"

# ── exact UPI PPT palette ──────────────────────────────────
NAVY    = RGBColor(0x1E,0x3A,0x5F)
TEAL    = RGBColor(0x2E,0x8B,0x7B)
BLUE    = RGBColor(0x3B,0x82,0xF6)
GREEN   = RGBColor(0x10,0xB9,0x81)
PURPLE  = RGBColor(0x8B,0x5C,0xF6)
RED     = RGBColor(0xEF,0x44,0x44)
AMBER   = RGBColor(0xF5,0x9E,0x0B)
GRAY    = RGBColor(0x6B,0x72,0x80)
DARK    = RGBColor(0x1E,0x29,0x3B)
WHITE   = RGBColor(0xFF,0xFF,0xFF)

C_BLUE  = RGBColor(0xDB,0xEA,0xFE)
C_GREEN = RGBColor(0xD1,0xFA,0xE5)
C_PURP  = RGBColor(0xED,0xE9,0xFE)
C_AMBER = RGBColor(0xFE,0xF3,0xC7)
C_RED   = RGBColor(0xFF,0xED,0xD5)
C_GRAY  = RGBColor(0xE5,0xE7,0xEB)

# Base Font
FONT_NAME = "Segoe UI"

blank = prs.slide_layouts[6]

def slide(): 
    sl = prs.slides.add_slide(blank)
    rect(sl, 0, 0, W, H, fill=WHITE, shape_type=MSO_SHAPE.RECTANGLE)
    return sl

def rect(sl, l,t,w,h, fill=None, line_col=None, line_w=Pt(0), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = sl.shapes.add_shape(shape_type,l,t,w,h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_col:
        s.line.color.rgb = line_col; s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def tb(sl, text, l,t,w,h, sz=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False):
    box = sl.shapes.add_textbox(l,t,w,h)
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = FONT_NAME
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return box

def hdr(sl, title):
    rect(sl, Inches(0.35), Inches(0.18), Inches(0.1), Inches(0.6), fill=TEAL, shape_type=MSO_SHAPE.RECTANGLE)
    tb(sl, title, Inches(0.55), Inches(0.14), Inches(12.2), Inches(0.7), sz=30, bold=True, color=NAVY)
    rect(sl, Inches(0.35), Inches(0.82), Inches(12.6), Pt(2), fill=TEAL, shape_type=MSO_SHAPE.RECTANGLE)

def card(sl, l,t,w,h, fill=C_BLUE, border=BLUE):
    return rect(sl,l,t,w,h, fill=fill, line_col=border, line_w=Pt(1), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

def img(sl, path, l,t,w,h=None):
    if os.path.exists(path):
        if h: sl.shapes.add_picture(path,l,t,w,h)
        else: sl.shapes.add_picture(path,l,t,w)

LOGO = r"d:\miniProject\Spam_reporting_portal\Report\iiit manipur.png"
PI   = r"d:\miniProject\Spam_reporting_portal\Report\ppt_imgs"

# ════════════════════════════════════════════════════════════
# TITLE SLIDE 
# ════════════════════════════════════════════════════════════
s0 = slide()
rect(s0, 0, 0, W, Inches(0.12), fill=TEAL, shape_type=MSO_SHAPE.RECTANGLE)
rect(s0, 0, H-Inches(0.12), W, Inches(0.12), fill=TEAL, shape_type=MSO_SHAPE.RECTANGLE)
rect(s0, Inches(2.5), Inches(2.22), Inches(8.3), Pt(2.5), fill=TEAL, shape_type=MSO_SHAPE.RECTANGLE)
rect(s0, Inches(2.5), Inches(2.32), Inches(8.3), Pt(1),   fill=C_GREEN, shape_type=MSO_SHAPE.RECTANGLE)

tb(s0, "CyberShield: Uncertainty-Aware Fraud Detection", Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.85), sz=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
tb(s0, "with Human-in-the-Loop Calibration for Cybercrime Intelligence", Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.6), sz=26, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
tb(s0, "A report submitted for the course Project – I (CS3201)", Inches(0.5), Inches(1.82), Inches(12.3), Inches(0.4), sz=16, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

tb(s0, "Submitted By",   Inches(1.8), Inches(2.55), Inches(4), Inches(0.38), sz=18, bold=True, color=DARK)
tb(s0, "ABHISHEK YADAV", Inches(1.8), Inches(2.95), Inches(4), Inches(0.4), sz=22, bold=True, color=NAVY)
tb(s0, "Roll No. 230103041  |  Semester VI", Inches(1.8), Inches(3.38), Inches(4), Inches(0.38), sz=16, color=GRAY)

rect(s0, Inches(6.55), Inches(2.55), Pt(1.5), Inches(1.3), fill=C_GRAY, shape_type=MSO_SHAPE.RECTANGLE)

tb(s0, "Supervised By", Inches(6.9), Inches(2.55), Inches(5), Inches(0.38), sz=18, bold=True, color=DARK)
tb(s0, "Dr. Salam Michael Singh", Inches(6.9), Inches(2.95), Inches(5), Inches(0.4), sz=22, bold=True, color=NAVY)

if os.path.exists(LOGO): img(s0, LOGO, Inches(5.55), Inches(3.95), Inches(2.2))

tb(s0,"Department of Computer Science and Engineering", Inches(0.5),Inches(6.35),Inches(12.3),Inches(0.35), sz=15,color=GRAY,align=PP_ALIGN.CENTER)
tb(s0,"Indian Institute of Information Technology Senapati, Manipur", Inches(0.5),Inches(6.7),Inches(12.3),Inches(0.35), sz=15,color=GRAY,align=PP_ALIGN.CENTER)
tb(s0,"April, 2026", Inches(0.5),Inches(7.05),Inches(12.3),Inches(0.35), sz=16,bold=True,color=TEAL,align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 1 – INTRODUCTION, OBJECTIVE & CONTRIBUTION
# ════════════════════════════════════════════════════════════
s1 = slide()
hdr(s1,"Introduction, Objective & Contribution")

# ── Box 1: THE PROBLEM (3 items) ─────────────────────────────
card(s1, Inches(0.3), Inches(1.0), Inches(5.8), Inches(3.4), C_BLUE, BLUE)
tb(s1,"1.  THE PROBLEM", Inches(0.55),Inches(1.1),Inches(5.3),Inches(0.4), sz=15,bold=True,color=BLUE)
probs=[
    ("Binary Rigidity in Existing Systems",
     "Standard ML forces a hard Fraud/Legit output, making arbitrary commitments on messages sitting at the decision boundary."),
    ("Concept Drift & Semantic Evasion",
     "Attackers rotate domains every 48 hours and exploit character obfuscation, rendering static models obsolete within days."),
    ("Fragmented Reporting Workflow",
     "Victims must manually identify IoCs and draft NCRP complaints from scratch, causing transcription errors and incomplete evidence."),
]
for i,(t,p) in enumerate(probs):
    rect(s1,Inches(0.55),Inches(1.62)+i*Inches(1.0), Inches(0.12),Inches(0.12),fill=BLUE, shape_type=MSO_SHAPE.OVAL)
    tb(s1,t, Inches(0.75),Inches(1.57)+i*Inches(1.0), Inches(5.1),Inches(0.3), sz=13,bold=True,color=NAVY)
    tb(s1,p, Inches(0.75),Inches(1.82)+i*Inches(1.0), Inches(5.1),Inches(0.65), sz=12,color=DARK)

# ── Box 2: PRIMARY OBJECTIVES (3 items) ──────────────────────
card(s1, Inches(0.3), Inches(4.5), Inches(5.8), Inches(2.75), C_AMBER, AMBER)
tb(s1,"2.  PRIMARY OBJECTIVES", Inches(0.55),Inches(4.62),Inches(5.3),Inches(0.4), sz=15,bold=True,color=AMBER)
objs=[
    "Build a probabilistic, uncertainty-aware ensemble (LR + SVM + NB) that classifies across 4 risk strata instead of forcing binary outcomes.",
    "Integrate a Human-in-the-Loop calibration loop with Hard Negative Mining to resist concept drift continuously.",
    "Generate automated NCRP-compliant PDF reports bridging AI predictions to law enforcement action.",
]
for i,p in enumerate(objs):
    rect(s1,Inches(0.55),Inches(5.13)+i*Inches(0.72), Inches(0.12),Inches(0.12),fill=AMBER, shape_type=MSO_SHAPE.OVAL)
    tb(s1,p, Inches(0.75),Inches(5.08)+i*Inches(0.72), Inches(5.1),Inches(0.62), sz=12,color=DARK)

# ── Box 3: 5 UNIQUE CONTRIBUTIONS (updated wording) ──────────
card(s1,Inches(6.3),Inches(1.0),Inches(6.7),Inches(6.25),C_GREEN,GREEN)
tb(s1,"3.  CYBERSHIELD: 5 UNIQUE CONTRIBUTIONS",Inches(6.55),Inches(1.1),Inches(6.2),Inches(0.4), sz=15,bold=True,color=GREEN)
contribs=[
    (BLUE,  "Uncertainty-Aware Classification",
     "4-tier probability routing (FRAUD \u22650.90 / SUSPICIOUS \u22650.70 / UNCERTAIN 0.30\u20130.70 / LEGIT <0.30) replacing forced binary predictions."),
    (GREEN, "Human-in-the-Loop Calibration",
     "UNCERTAIN cases routed to admin panel; human labels feed Hard Negative Mining retraining pipeline automatically."),
    (PURPLE,"Token-Level Explainability",
     "Word Risk Heatmap maps Logistic Regression coefficient weights to individual tokens, showing victims why a message was flagged."),
    (AMBER, "Automated NCRP Adjudication",
     "One-click PDF generation embedding complaint ID, extracted IoCs, risk level, and NCRP category code \u2014 ready for submission."),
    (RED,   "Velocity Intelligence Tracking",
     "Cross-submission frequency tracker escalating repeated phone/URL artifacts: MED (1 hit) \u2192 HIGH (\u22652) \u2192 CRITICAL (\u22654)."),
]
for i,(col,title,desc) in enumerate(contribs):
    y=Inches(1.65)+i*Inches(1.1)
    dot=s1.shapes.add_shape(MSO_SHAPE.OVAL,Inches(6.55),y+Inches(0.08), Inches(0.15),Inches(0.15))
    dot.fill.solid();dot.fill.fore_color.rgb=col;dot.line.fill.background()
    tb(s1,title,Inches(6.8),y,         Inches(6.1),Inches(0.32), sz=14,bold=True,color=NAVY)
    tb(s1,desc, Inches(6.8),y+Inches(0.3),Inches(6.1),Inches(0.68), sz=12,color=DARK)

# ════════════════════════════════════════════════════════════
# SLIDE 2 – BACKGROUND / EXISTING WORK  (dark text-table style)
# ════════════════════════════════════════════════════════════
s2=slide()
hdr(s2,"Background / Existing Work")

# ── Colours specific to this table ─────────────────────────
DIVIDER  = RGBColor(0x3B,0x4A,0x5E)   # subtle separator line
X_COL    = RGBColor(0x94,0xA3,0xB8)   # muted grey for X cells
CHK_COL  = GREEN                       # green for ✓ CyberShield cells

# ── Column layout ──────────────────────────────────────────
COL_HDRS = ["System", "Algorithm", "Uncertainty\nHandling", "HITL",
            "Artifact\nExtraction", "Reporting"]
col_x = [Inches(0.45), Inches(2.55), Inches(4.65), Inches(6.95),
         Inches(8.75), Inches(10.95)]
col_w = [Inches(2.05), Inches(2.05), Inches(2.25), Inches(1.75),
         Inches(2.15), Inches(1.95)]

# ── Row data: (System, Algorithm, Uncertainty, HITL, Artifact, Reporting)
ROWS2 = [
    ("Almeida et al.\n(2011)", "SVM, Naive Bayes",
     "None \u2014 Binary\nonly",           "X", "X",          "X"),
    ("Roy et al.\n(2020)",     "CNN / LSTM",
     "None \u2014 Black-\nbox",            "X", "X",          "X"),
    ("Jain & Gupta\n(2021)",   "Lexical URL + NLP",
     "None",                               "X", "URLs only",  "X"),
    ("Abid et al.\n(2022)",    "Optimized LR/SVM",
     "None",                               "X", "X",          "X"),
    ("CyberShield\n(2026)",    "Ensemble\nLR+SVM+NB",
     "4-Tier\nProbabilistic",  "\u2713 Active\nLearning",
     "URLs +\nPhones",         "\u2713 NCRP\nPDF"),
]

TABLE_TOP2 = Inches(1.05)
HDR_H2     = Inches(0.70)
ROW_H2     = Inches(1.12)
TABLE_W2   = Inches(12.73)
TABLE_H2   = HDR_H2 + len(ROWS2) * ROW_H2

# Light theme card — NAVY border
rect(s2, Inches(0.3), TABLE_TOP2, TABLE_W2, TABLE_H2,
     fill=WHITE, line_col=NAVY, line_w=Pt(2.5), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

# Header bar — NAVY fill
rect(s2, Inches(0.3), TABLE_TOP2, TABLE_W2, HDR_H2,
     fill=NAVY, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s2, Inches(0.3), TABLE_TOP2+HDR_H2-Inches(0.18), TABLE_W2, Inches(0.18),
     fill=NAVY, shape_type=MSO_SHAPE.RECTANGLE)   # flatten header bottom

# Column headers — white bold text
for j,(hdr_txt,cx,cw) in enumerate(zip(COL_HDRS,col_x,col_w)):
    tb(s2, hdr_txt, cx, TABLE_TOP2+Inches(0.08), cw-Inches(0.05), HDR_H2-Inches(0.1),
       sz=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# ── Data rows ──────────────────────────────────────────────
for i, row in enumerate(ROWS2):
    sys_name = row[0]
    cells    = row[1:]
    ry       = TABLE_TOP2 + HDR_H2 + i * ROW_H2
    is_cs    = "CyberShield" in sys_name

    # Alternating row background
    if i % 2 == 1:
        rect(s2, Inches(0.31), ry, Inches(12.71), ROW_H2,
             fill=C_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

    # System name
    tb(s2, sys_name, col_x[0], ry+Inches(0.15),
       col_w[0]-Inches(0.1), ROW_H2-Inches(0.2),
       sz=13, bold=is_cs, color=DARK)

    # Other cells
    for j,(cell,cx,cw) in enumerate(zip(cells, col_x[1:], col_w[1:])):
        is_check = cell.startswith("\u2713")
        is_x     = cell.strip() == "X"
        
        c_col = TEAL if is_cs else DARK
        c_bold = True if is_cs else False

        tb(s2, cell, cx, ry+Inches(0.15), cw-Inches(0.1), ROW_H2-Inches(0.2),
           sz=13, bold=c_bold, color=c_col)






# ════════════════════════════════════════════════════════════
# SLIDE 3 – PROPOSED SYSTEM / ARCHITECTURE
# ════════════════════════════════════════════════════════════
s3=slide()
hdr(s3,"Proposed System / Architecture")

pipe=os.path.join(PI,"pipeline_flow.png")
if os.path.exists(pipe): 
    img(s3,pipe,Inches(0.3),Inches(1.05),Inches(7.2),Inches(5.2))

card(s3, Inches(0.3), Inches(6.4), Inches(7.2), Inches(0.8), C_GRAY, DARK)
flow_txt = "END-TO-END FLOW: Citizen Input \u2192 RegEx & TF-IDF \u2192 Ensemble Scoring \u2192 4-Tier Routing \u2192 Velocity DB \u2192 HITL Queue or NCRP PDF"
tb(s3, flow_txt, Inches(0.4), Inches(6.5), Inches(7.0), Inches(0.6), sz=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)

# Top Right: ML Workflow Card
card(s3,Inches(7.7),Inches(1.05),Inches(5.3),Inches(2.8),C_AMBER,AMBER)
tb(s3,"ML WORKFLOW",Inches(7.85),Inches(1.15),Inches(5.0),Inches(0.35), sz=14,bold=True,color=AMBER)

ml_pipe=os.path.join(PI,"ml_pipeline.png")
if os.path.exists(ml_pipe): 
    img(s3,ml_pipe,Inches(7.8),Inches(1.5),Inches(5.1),Inches(2.2))

# Bottom Right: 4-Tier Routing Card
card(s3,Inches(7.7),Inches(3.95),Inches(5.3),Inches(3.1),C_BLUE,BLUE)
tb(s3,"4-TIER PROBABILISTIC ROUTING",Inches(7.85),Inches(4.05),Inches(5.0),Inches(0.35), sz=14,bold=True,color=BLUE)

thr_logic_img=os.path.join(PI,"routing_flow.png")
if os.path.exists(thr_logic_img): 
    img(s3,thr_logic_img,Inches(7.8),Inches(4.45),Inches(5.1),Inches(2.5))


# ════════════════════════════════════════════════════════════
# SLIDE 4 – RESULT & ANALYSIS  (new layout)
# ════════════════════════════════════════════════════════════
s4=slide()
hdr(s4,"Result & Analysis")

# ── LEFT column (Graphs)  x=0.3  w=6.3 ────────────────────────
LX = Inches(0.3); LW = Inches(6.3)
roc=os.path.join(PI,"fig_roc_pr_curves.png")
if os.path.exists(roc): 
    rect(s4, LX-Inches(0.02), Inches(1.2)-Inches(0.02), LW+Inches(0.04), Inches(3.0)+Inches(0.04), fill=WHITE, line_col=NAVY, line_w=Pt(2.5), shape_type=MSO_SHAPE.RECTANGLE)
    img(s4,roc,LX,Inches(1.2),LW,Inches(3.0))

thr=os.path.join(PI,"fig_threshold_sensitivity.png")
if os.path.exists(thr): 
    rect(s4, LX-Inches(0.02), Inches(4.4)-Inches(0.02), LW+Inches(0.04), Inches(2.8)+Inches(0.04), fill=WHITE, line_col=NAVY, line_w=Pt(2.5), shape_type=MSO_SHAPE.RECTANGLE)
    img(s4,thr,LX,Inches(4.4),LW,Inches(2.8))

# ── RIGHT column (Tables + CM)  x=6.8  w=6.3 ──────────────────
RX = Inches(6.8); RW = Inches(6.3)

# Performance Table
PT_TOP = Inches(1.2); PT_HDR = Inches(0.4); PT_ROW = Inches(0.35)
PT_ROWS = [
    ("Logistic Regression","0.9324","0.9512","0.9417","0.9871"),
    ("SVM (Calibrated)",   "0.9473","0.9419","0.9446","0.9892"),
    ("Naive Bayes",        "0.9405","0.9346","0.9376","0.9864"),
    ("Ensemble \u2605",    "0.9518","0.9419","0.9467","0.9895"),
]
PT_COLS = ["Model","Precision","Recall","F1","ROC-AUC"]
pt_cw   = [Inches(2.1),Inches(1.0),Inches(0.9),Inches(0.9),Inches(1.2)]
pt_cx   = [RX+Inches(0.05)]
for w in pt_cw[:-1]: pt_cx.append(pt_cx[-1]+w)
PT_H = PT_HDR + len(PT_ROWS)*PT_ROW

tb(s4,"Individual Model vs Ensemble Performance",RX,PT_TOP-Inches(0.28),RW,Inches(0.25),
   sz=11,bold=True,color=NAVY)
rect(s4,RX,PT_TOP,RW,PT_H, fill=WHITE, line_col=NAVY, line_w=Pt(2.5), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s4,RX,PT_TOP,RW,PT_HDR, fill=NAVY, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s4,RX,PT_TOP+PT_HDR-Inches(0.18),RW,Inches(0.18), fill=NAVY, shape_type=MSO_SHAPE.RECTANGLE)

for j,(h,cx,cw) in enumerate(zip(PT_COLS,pt_cx,pt_cw)):
    tb(s4,h,cx,PT_TOP,cw-Inches(0.04),PT_HDR,
       sz=11,bold=True,color=WHITE,align=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)

for i,(row) in enumerate(PT_ROWS):
    ry=PT_TOP+PT_HDR+i*PT_ROW
    is_ens="Ensemble" in row[0]
    if i % 2 == 1:
        rect(s4,RX+Inches(0.02),ry,RW-Inches(0.04),PT_ROW, fill=C_BLUE, shape_type=MSO_SHAPE.RECTANGLE)
    for j,(cell,cx,cw) in enumerate(zip(row,pt_cx,pt_cw)):
        c=TEAL if is_ens else DARK
        tb(s4,cell,cx,ry+Inches(0.02),cw-Inches(0.04),PT_ROW,
           sz=10,bold=is_ens,color=c,align=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)

# HNM Ablation table
AB_TOP = Inches(3.2); AB_HDR = Inches(0.4); AB_ROW = Inches(0.35)
AB_ROWS=[
    ("Accuracy",      "0.9482","0.9487","+0.0005"),
    ("Precision",     "0.9508","0.9518","+0.0010"),
    ("Recall (Fraud)","0.9419","0.9419"," 0.0000"),
    ("Macro F1",      "0.9482","0.9486","+0.0004"),
    ("ROC-AUC",       "0.9893","0.9895","+0.0002"),
]
AB_COLS=["Metric","Pre-HNM","Post-HNM","\u0394"]
ab_cw=[Inches(1.8),Inches(1.2),Inches(1.2),Inches(1.2)]
ab_cx=[RX+Inches(0.05)]
for w in ab_cw[:-1]: ab_cx.append(ab_cx[-1]+w)
AB_H=AB_HDR+len(AB_ROWS)*AB_ROW

tb(s4,"FP-Driven HNM Ablation (Pre vs Post)",RX,AB_TOP-Inches(0.25),RW,Inches(0.22),
   sz=11,bold=True,color=NAVY)
rect(s4,RX,AB_TOP,RW,AB_H, fill=WHITE, line_col=NAVY, line_w=Pt(2.5), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s4,RX,AB_TOP,RW,AB_HDR, fill=NAVY, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s4,RX,AB_TOP+AB_HDR-Inches(0.18),RW,Inches(0.18), fill=NAVY, shape_type=MSO_SHAPE.RECTANGLE)

for j,(h,cx,cw) in enumerate(zip(AB_COLS,ab_cx,ab_cw)):
    tb(s4,h,cx,AB_TOP,cw-Inches(0.04),AB_HDR,
       sz=11,bold=True,color=WHITE,align=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)

for i,row in enumerate(AB_ROWS):
    ry=AB_TOP+AB_HDR+i*AB_ROW
    if i % 2 == 1:
        rect(s4,RX+Inches(0.02),ry,RW-Inches(0.04),AB_ROW, fill=C_BLUE, shape_type=MSO_SHAPE.RECTANGLE)
    for j,(cell,cx,cw) in enumerate(zip(row,ab_cx,ab_cw)):
        delta_green=j==3 and cell.startswith("+")
        c=GREEN if delta_green else(C_GRAY if(j==3 and cell.strip()=="0.0000") else DARK)
        tb(s4,cell,cx,ry+Inches(0.02),cw-Inches(0.04),AB_ROW,
           sz=10,color=c,align=PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)

# Confusion Matrix
cm=os.path.join(PI,"fig_confusion_matrices.png")
if os.path.exists(cm): 
    rect(s4, RX-Inches(0.02), Inches(5.5)-Inches(0.02), RW+Inches(0.04), Inches(1.7)+Inches(0.04), fill=WHITE, line_col=NAVY, line_w=Pt(2.5), shape_type=MSO_SHAPE.RECTANGLE)
    img(s4,cm,RX,Inches(5.5),RW,Inches(1.7))



# ════════════════════════════════════════════════════════════
# SLIDE 5 – CONCLUSION
# ════════════════════════════════════════════════════════════
s5=slide()
hdr(s5,"Conclusion & Future Scope")

card(s5, Inches(0.3), Inches(1.0), Inches(7.5), Inches(5.9), C_BLUE, BLUE)
tb(s5,"SUMMARY OF WORK DONE", Inches(0.5),Inches(1.1),Inches(7.0),Inches(0.4), sz=15,bold=True,color=BLUE)

summ=[
    ("Uncertainty-Aware ML Ensemble", "The calibrated soft-voting ensemble (LR + SVM + NB) over a 25,000-feature TF-IDF matrix achieved 94.70% test-set accuracy and ROC-AUC of 0.9900 on 3,980 frozen adversarial samples \u2014 outperforming every individual base classifier."),
    ("System Safety via the UNCERTAIN Band", "The wide UNCERTAIN zone (0.30 \u2264 P < 0.70) is not a limitation \u2014 it is a deliberate safety mechanism that prevents low-confidence binary commitments, routing ambiguous cases to human experts rather than guessing. No existing system implements this."),
    ("Bridging Predictions to Law Enforcement Action", "By combining Velocity Intelligence tracking with automated NCRP-compliant PDF generation, CyberShield eliminates the gap between a raw model probability and a ready-to-submit cybercrime complaint \u2014 a gap every prior academic system ignored."),
    ("Adaptive by Design", "Hard Negative Mining + HITL retraining ensures the model self-calibrates against new threat vocabulary without engineering overhead, reducing fraud-class False Negatives from 112 \u2192 108 in a single cycle.")
]
for i,(label,desc) in enumerate(summ):
    y=Inches(1.6)+i*Inches(1.05)
    dot=s5.shapes.add_shape(MSO_SHAPE.OVAL,Inches(0.5),y+Inches(0.06), Inches(0.12),Inches(0.12))
    dot.fill.solid();dot.fill.fore_color.rgb=NAVY;dot.line.fill.background()
    tb(s5,label,Inches(0.7),y,Inches(6.8),Inches(0.3), sz=14,bold=True,color=NAVY)
    tb(s5,desc, Inches(0.7),y+Inches(0.25),Inches(6.8),Inches(0.85), sz=12,color=DARK) # Adjusted size slightly for longer text

card(s5,Inches(8.0),Inches(1.0),Inches(5.0),Inches(5.9),C_PURP,PURPLE)
tb(s5,"FUTURE SCOPE",Inches(8.2),Inches(1.1),Inches(4.6),Inches(0.4), sz=15,bold=True,color=PURPLE)
future=[
    (GREEN, "Near-Term", "Direct REST API integration with cybercrime.gov.in to auto-submit high-confidence (P \u2265 0.90) cases as draft FIRs, replacing the current PDF-and-helper model."),
    (GREEN, "Near-Term", "Browser extension packaging the inference endpoint for in-context WhatsApp/SMS analysis without visiting the portal."),
    (AMBER, "Mid-Term", "MuRIL / IndicBERT multilingual pipeline supporting Devanagari, Meitei Mayek, and all 22 scheduled Indian languages."),
    (AMBER, "Mid-Term", "DistilBERT contextual embeddings replacing TF-IDF, contingent on meeting the 500ms inference latency budget on CPU infrastructure."),
    (RED,   "Long-Term", "Federated Learning architecture enabling decentralized model updates across state cybercrime offices without exposing victim message data."),
]
for i,(col,tag,desc) in enumerate(future):
    y=Inches(1.6)+i*Inches(0.9) # Adjusted gap slightly
    rect(s5,Inches(8.2),y,Inches(1.1),Inches(0.28),fill=col, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s5,tag,Inches(8.2),y-Inches(0.02),Inches(1.1),Inches(0.3), sz=10,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    tb(s5,desc,Inches(9.4),y-Inches(0.1),Inches(3.4),Inches(1.0),sz=12,color=DARK) # Adjusted font to fit larger descriptions

# ════════════════════════════════════════════════════════════
# SLIDE 6 – REFERENCES
# ════════════════════════════════════════════════════════════
s6=slide()
hdr(s6,"References")
refs=[
    "[1] Almeida, T.A., Hidalgo, J.M.G., & Yamakami, A. (2011). Contributions to the study of SMS spam filtering. ACM DocEng, 259-262.",
    "[2] Settles, B. (2012). Active Learning. Synthesis Lectures on AI and Machine Learning. Morgan & Claypool.",
    "[3] Holzinger, A. (2016). Interactive machine learning for health informatics. Brain Informatics, 3(2), 119-131.",
    "[4] Chawla, N.V. et al. (2002). SMOTE: Synthetic Minority Over-sampling Technique. JAIR, 16, 321-357.",
    "[5] Roy, S. et al. (2020). A Deep Learning based AI Model for SMS Spam Detection. IJEAT, 9(3).",
    "[6] Pedregosa, F. et al. (2011). Scikit-learn: Machine Learning in Python. JMLR, 12, 2825-2830.",
    "[7] Purkait, S. (2014). Phishing counter measures and their effectiveness. Info. Mgmt. & Comp. Security, 22(5).",
    "[8] Ministry of Home Affairs, Govt. of India (2023). Annual Report on Cybercrime - NCRP Statistics 2023.",
]
for i,ref in enumerate(refs):
    rect(s6,Inches(0.3),Inches(1.02)+i*Inches(0.7), Inches(12.6),Inches(0.6), fill=C_BLUE if i%2==0 else WHITE, line_col=DARK, line_w=Pt(1.5), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s6,ref,Inches(0.45),Inches(1.15)+i*Inches(0.7), Inches(12.3),Inches(0.5),sz=13,color=DARK)

# ════════════════════════════════════════════════════════════
# SLIDE 7 – THANK YOU
# ════════════════════════════════════════════════════════════
s7 = slide()
rect(s7, 0, 0, W, Inches(0.12), fill=TEAL, shape_type=MSO_SHAPE.RECTANGLE)
rect(s7, 0, H-Inches(0.12), W, Inches(0.12), fill=TEAL, shape_type=MSO_SHAPE.RECTANGLE)

tb(s7, "Thank You", Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5), sz=70, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
rect(s7, Inches(5.66), Inches(4.0), Inches(2.0), Pt(2), fill=TEAL, shape_type=MSO_SHAPE.RECTANGLE)
tb(s7, "Open to Questions & Discussion", Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6), sz=24, bold=False, color=DARK, align=PP_ALIGN.CENTER)

out = r"d:\miniProject\Spam_reporting_portal\CyberShield_Presentation_v40.pptx"
prs.save(out)
print("Saved:", out)
