
# CyberShield PPT – styled after Smart UPI Fraud Detection.pptx palette
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

# ── exact UPI PPT palette ──────────────────────────────────
BG      = RGBColor(0xF9,0xFA,0xFB)  # slide background (near white)
NAVY    = RGBColor(0x1E,0x3A,0x5F)  # heading text
TEAL    = RGBColor(0x2E,0x8B,0x7B)  # accent / bars
BLUE    = RGBColor(0x3B,0x82,0xF6)  # highlight
GREEN   = RGBColor(0x10,0xB9,0x81)  # success
PURPLE  = RGBColor(0x8B,0x5C,0xF6)  # special
RED     = RGBColor(0xEF,0x44,0x44)  # warning
AMBER   = RGBColor(0xF5,0x9E,0x0B)  # caution
GRAY    = RGBColor(0x6B,0x72,0x80)  # subtext
DARK    = RGBColor(0x1E,0x29,0x3B)  # dark text
WHITE   = RGBColor(0xFF,0xFF,0xFF)
# card fills
C_BLUE  = RGBColor(0xDB,0xEA,0xFE)
C_GREEN = RGBColor(0xD1,0xFA,0xE5)
C_PURP  = RGBColor(0xED,0xE9,0xFE)
C_AMBER = RGBColor(0xFE,0xF3,0xC7)
C_RED   = RGBColor(0xFF,0xED,0xD5)
C_GRAY  = RGBColor(0xE5,0xE7,0xEB)

blank = prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(blank)

def rect(sl, l,t,w,h, fill=BG, line_col=None, line_w=Pt(0)):
    s = sl.shapes.add_shape(1,l,t,w,h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_col:
        s.line.color.rgb = line_col; s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def tb(sl, text, l,t,w,h, sz=13, bold=False, color=DARK,
       align=PP_ALIGN.LEFT, italic=False):
    box = sl.shapes.add_textbox(l,t,w,h)
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return box

def hdr(sl, title):
    # teal left bar
    rect(sl, Inches(0.35), Inches(0.18), Inches(0.1), Inches(0.6), fill=TEAL)
    tb(sl, title, Inches(0.55), Inches(0.14), Inches(12.2), Inches(0.7),
       sz=28, bold=True, color=NAVY)
    # underline
    rect(sl, Inches(0.35), Inches(0.82), Inches(12.6), Pt(2), fill=TEAL)

def card(sl, l,t,w,h, fill=C_BLUE, border=BLUE):
    return rect(sl,l,t,w,h, fill=fill, line_col=border, line_w=Pt(1))

def img(sl, path, l,t,w,h=None):
    if os.path.exists(path):
        if h: sl.shapes.add_picture(path,l,t,w,h)
        else: sl.shapes.add_picture(path,l,t,w)

LOGO = r"d:\miniProject\Spam_reporting_portal\Report\iiit manipur.png"
PI   = r"d:\miniProject\Spam_reporting_portal\Report\ppt_imgs"

# ════════════════════════════════════════════════════════════
# TITLE SLIDE  – white, matching UPI format exactly
# ════════════════════════════════════════════════════════════
s0 = slide()
rect(s0,0,0,W,H, fill=WHITE)
# top teal bar
rect(s0, 0, 0, W, Inches(0.12), fill=TEAL)
# bottom teal bar
rect(s0, 0, H-Inches(0.12), W, Inches(0.12), fill=TEAL)
# two teal accent lines under title
rect(s0, Inches(2.5), Inches(2.22), Inches(8.3), Pt(2.5), fill=TEAL)
rect(s0, Inches(2.5), Inches(2.32), Inches(8.3), Pt(1),   fill=C_GREEN)

tb(s0, "CyberShield: Uncertainty-Aware Fraud Detection",
   Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.85),
   sz=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
tb(s0, "with Human-in-the-Loop Calibration for Cybercrime Intelligence",
   Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.6),
   sz=22, color=TEAL, align=PP_ALIGN.CENTER)
tb(s0, "A report submitted for the course Project – I (CS3201)",
   Inches(0.5), Inches(1.82), Inches(12.3), Inches(0.4),
   sz=14, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

# Submitted / supervised cols
tb(s0, "Submitted By",   Inches(1.8), Inches(2.55), Inches(4), Inches(0.38),
   sz=16, bold=True, color=DARK)
tb(s0, "ABHISHEK YADAV", Inches(1.8), Inches(2.95), Inches(4), Inches(0.4),
   sz=19, bold=True, color=NAVY)
tb(s0, "Roll No. 230103041  |  Semester VI",
   Inches(1.8), Inches(3.38), Inches(4), Inches(0.38), sz=14, color=GRAY)

rect(s0, Inches(6.55), Inches(2.55), Pt(1.5), Inches(1.3), fill=C_GRAY)

tb(s0, "Supervised By",       Inches(6.9), Inches(2.55), Inches(5), Inches(0.38),
   sz=16, bold=True, color=DARK)
tb(s0, "Dr. Salam Michael Singh", Inches(6.9), Inches(2.95), Inches(5), Inches(0.4),
   sz=19, bold=True, color=NAVY)

if os.path.exists(LOGO):
    img(s0, LOGO, Inches(5.55), Inches(3.95), Inches(2.2))

tb(s0,"Department of Computer Science and Engineering",
   Inches(0.5),Inches(6.35),Inches(12.3),Inches(0.35),
   sz=13,color=GRAY,align=PP_ALIGN.CENTER)
tb(s0,"Indian Institute of Information Technology Senapati, Manipur",
   Inches(0.5),Inches(6.7),Inches(12.3),Inches(0.35),
   sz=13,color=GRAY,align=PP_ALIGN.CENTER)
tb(s0,"April, 2026",
   Inches(0.5),Inches(7.05),Inches(12.3),Inches(0.35),
   sz=13,bold=True,color=TEAL,align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 1 – INTRODUCTION, OBJECTIVES & CONTRIBUTIONS
# ════════════════════════════════════════════════════════════
s1 = slide()
rect(s1,0,0,W,H,fill=BG)
hdr(s1,"Introduction, Objectives & Unique Contributions")

# LEFT – Problem panel
card(s1, Inches(0.3), Inches(1.0), Inches(4.1), Inches(6.1), C_BLUE, BLUE)
tb(s1,"THE PROBLEM", Inches(0.45),Inches(1.08),Inches(3.8),Inches(0.42),
   sz=13,bold=True,color=BLUE)
probs=[
    ("India recorded 1.12M+ cybercrime complaints in 2023",NAVY),
    ("Binary spam filters — no uncertainty, no nuance",DARK),
    ("ML models drift as scammers evolve language",DARK),
    ("Victims have no structured cyber-complaint path",DARK),
    ("No visibility into repeat malicious URLs/phones",DARK),
]
for i,(p,c) in enumerate(probs):
    rect(s1,Inches(0.45),Inches(1.58)+i*Inches(0.98),
         Inches(0.12),Inches(0.12),fill=BLUE)
    tb(s1,p,Inches(0.65),Inches(1.53)+i*Inches(0.98),
       Inches(3.65),Inches(0.82),sz=12,color=c)

# RIGHT – Contributions
card(s1,Inches(4.6),Inches(1.0),Inches(8.4),Inches(6.1),C_GREEN,GREEN)
tb(s1,"UNIQUE CONTRIBUTIONS",Inches(4.75),Inches(1.08),Inches(8.1),Inches(0.42),
   sz=13,bold=True,color=GREEN)
contribs=[
    (BLUE,  "4-Tier Probabilistic Thresholding",
             "FRAUD / SUSPICIOUS / UNCERTAIN / LEGIT strata replace binary\noutput — first of its kind for SMS fraud classification"),
    (GREEN, "Human-in-the-Loop Retraining Loop",
             "UNCERTAIN cases routed to admin queue, then Hard Negative\nMining re-calibrates model — automatic concept drift resistance"),
    (PURPLE,"Automated NCRP PDF Generation",
             "Auto-fills Indian cybercrime complaint with IoCs, probability\nscore & ID — zero-friction victim reporting"),
    (AMBER, "Velocity Intelligence Tracker",
             "Real-time frequency DB for malicious URLs & phones;\nMED/HIGH/CRITICAL threat tiers across all submissions"),
    (RED,   "Token-Level Word Risk Heatmap",
             "LR coefficient weights surface fraud-driving tokens per\nmessage — explainable AI for non-technical citizens"),
]
card_cols=[C_BLUE,C_GREEN,C_PURP,C_AMBER,C_RED]
for i,(col,title,desc) in enumerate(contribs):
    y=Inches(1.58)+i*Inches(1.04)
    card(s1,Inches(4.75),y,Inches(8.0),Inches(0.95),card_cols[i],col)
    dot=s1.shapes.add_shape(9,Inches(4.88),y+Inches(0.25),
                             Inches(0.22),Inches(0.22))
    dot.fill.solid();dot.fill.fore_color.rgb=col;dot.line.fill.background()
    tb(s1,title,Inches(5.18),y+Inches(0.06),Inches(7.4),Inches(0.38),
       sz=13,bold=True,color=col)
    tb(s1,desc, Inches(5.18),y+Inches(0.44),Inches(7.4),Inches(0.48),
       sz=11,color=GRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 2 – BACKGROUND / EXISTING WORK
# ════════════════════════════════════════════════════════════
s2=slide()
rect(s2,0,0,W,H,fill=BG)
hdr(s2,"Background & Existing Work")

headers=["System","Approach","Limitation","CyberShield Fix"]
rows=[
    ["SpamAssassin\n(Apache 2002)",   "Rule-based +\nBayesian Naive",  "Static rules break\non novel SMS slang",     "Ensemble + HNM\nretraining loop"],
    ["UCI SMS Spam\nBaseline",        "Single Naive Bayes\nTF-IDF",    "Binary only;\nno uncertainty band",          "4-tier probabilistic\nrouting"],
    ["MobiFish\n(Purkait 2014)",      "URL heuristics\nfor smishing",  "URL-only; ignores\nmessage semantics",        "NLP + artifact\ntokenization"],
    ["Deep Spam\n(Roy et al. 2020)",  "LSTM on word\nembeddings",      "Opaque; no citizen-\nfacing actionability",   "Word Risk Heatmap\n+ NCRP PDF"],
    ["NCRP Portal\n(Govt. India)",    "Manual complaint\nentry",        "High friction;\nno automated IoC",            "Pre-filled PDF +\nVelocity Tracker"],
]
cws=[Inches(2.5),Inches(2.5),Inches(3.0),Inches(3.0)]
cxs=[Inches(0.35),Inches(2.9),Inches(5.45),Inches(8.5)]
h_fills=[NAVY,BLUE,RED,GREEN]
for j,(h,cx,cw,hf) in enumerate(zip(headers,cxs,cws,h_fills)):
    rect(s2,cx,Inches(1.05),cw-Inches(0.06),Inches(0.48),fill=hf)
    tb(s2,h,cx+Inches(0.06),Inches(1.08),cw-Inches(0.1),Inches(0.42),
       sz=12,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
row_bg=[C_BLUE,BG,C_GREEN,BG,C_AMBER]
for i,row in enumerate(rows):
    y=Inches(1.57)+i*Inches(0.97)
    for j,(cell,cx,cw) in enumerate(zip(row,cxs,cws)):
        rect(s2,cx,y,cw-Inches(0.06),Inches(0.9),fill=row_bg[i],
             line_col=C_GRAY,line_w=Pt(0.5))
        tc=GREEN if j==3 else (RED if j==2 else DARK)
        tb(s2,cell,cx+Inches(0.08),y+Inches(0.08),cw-Inches(0.14),
           Inches(0.78),sz=11,color=tc)

# ════════════════════════════════════════════════════════════
# SLIDE 3 – PROPOSED SYSTEM / ARCHITECTURE
# ════════════════════════════════════════════════════════════
s3=slide()
rect(s3,0,0,W,H,fill=BG)
hdr(s3,"Proposed System — CyberShield Architecture")

pipe=os.path.join(PI,"pipeline_flow.png")
if os.path.exists(pipe):
    img(s3,pipe,Inches(0.3),Inches(1.0),Inches(7.0),Inches(5.8))

# right panel
card(s3,Inches(7.5),Inches(1.0),Inches(5.5),Inches(5.8),C_BLUE,BLUE)
tb(s3,"TECH STACK",Inches(7.65),Inches(1.08),Inches(5.2),Inches(0.4),
   sz=13,bold=True,color=BLUE)
stack=[
    (NAVY,  "Backend",       "Python 3.11 · Flask · SQLAlchemy · SQLite"),
    (BLUE,  "ML Engine",     "scikit-learn: LR + SVM + Naive Bayes\n25,000-feature FeatureUnion TF-IDF"),
    (GREEN, "Auth",          "Google OAuth 2.0 — admin route gating"),
    (PURPLE,"PDF Output",    "ReportLab — NCRP-compliant auto-fill"),
    (AMBER, "Explainability","LR coefficient weights → Word Risk Heatmap"),
    (TEAL,  "Deployment",    "Flask dev server · SQLite WAL mode"),
]
for i,(col,label,val) in enumerate(stack):
    y=Inches(1.55)+i*Inches(0.78)
    dot2=s3.shapes.add_shape(9,Inches(7.65),y+Inches(0.18),
                              Inches(0.16),Inches(0.16))
    dot2.fill.solid();dot2.fill.fore_color.rgb=col;dot2.line.fill.background()
    tb(s3,label,Inches(7.88),y,Inches(1.6),Inches(0.35),sz=11,bold=True,color=col)
    tb(s3,val,  Inches(9.5), y,Inches(3.35),Inches(0.65),sz=11,color=GRAY)

rect(s3,Inches(7.5),Inches(5.55),Inches(5.5),Pt(1.5),fill=TEAL)
tb(s3,"10-STAGE PIPELINE",Inches(7.65),Inches(5.65),Inches(5.2),Inches(0.35),
   sz=11,bold=True,color=NAVY)
tb(s3,"Input > Preprocess > TF-IDF > Ensemble Score > Threshold\nRouter > DB Log > Velocity Update > HITL Queue > NCRP PDF > Response",
   Inches(7.65),Inches(6.0),Inches(5.2),Inches(0.7),sz=10,color=GRAY)

# ════════════════════════════════════════════════════════════
# SLIDE 4 – RESULTS & ANALYSIS
# ════════════════════════════════════════════════════════════
s4=slide()
rect(s4,0,0,W,H,fill=BG)
hdr(s4,"Results & Analysis")

# KPI cards
kpis=[
    ("94.85%","Test Accuracy",   C_GREEN,GREEN),
    ("0.9899", "ROC-AUC",        C_BLUE, BLUE),
    ("26,534", "Training Samples",C_PURP,PURPLE),
    ("<500ms",  "Inference Speed",C_AMBER,AMBER),
    ("3,980",  "Frozen Test Set", C_RED, RED),
]
kw=Inches(2.42)
for i,(val,lab,bg2,col2) in enumerate(kpis):
    x=Inches(0.25)+i*(kw+Inches(0.06))
    card(s4,x,Inches(1.0),kw,Inches(1.2),bg2,col2)
    tb(s4,val,x,Inches(1.08),kw,Inches(0.58),sz=24,bold=True,
       color=col2,align=PP_ALIGN.CENTER)
    tb(s4,lab,x,Inches(1.65),kw,Inches(0.42),sz=11,color=GRAY,
       align=PP_ALIGN.CENTER)

# HNM ablation table
card(s4,Inches(0.25),Inches(2.35),Inches(4.6),Inches(4.5),C_GREEN,GREEN)
tb(s4,"HNM Ablation Table",Inches(0.4),Inches(2.42),Inches(4.3),Inches(0.38),
   sz=12,bold=True,color=GREEN)
abl_h=["Metric","Pre-HNM","Post-HNM","Delta"]
abl_r=[
    ["Accuracy",      "0.9482","0.9485","+0.0003"],
    ["Precision",     "0.9508","0.9494","-0.0014"],
    ["Recall(Fraud)", "0.9419","0.9440","+0.0021"],
    ["Macro F1",      "0.9482","0.9484","+0.0002"],
    ["ROC-AUC",       "0.9893","0.9899","+0.0006"],
]
cws2=[Inches(1.55),Inches(0.88),Inches(1.0),Inches(0.88)]
cxs2=[Inches(0.3),Inches(1.9),Inches(2.82),Inches(3.86)]
for j,(h,cx,cw) in enumerate(zip(abl_h,cxs2,cws2)):
    hf2=[NAVY,BLUE,GREEN,TEAL][j]
    rect(s4,cx,Inches(2.85),cw-Inches(0.04),Inches(0.38),fill=hf2)
    tb(s4,h,cx+Inches(0.04),Inches(2.87),cw,Inches(0.32),
       sz=10,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
for i,row in enumerate(abl_r):
    ry=Inches(3.26)+i*Inches(0.62)
    rbg=C_BLUE if i%2==0 else WHITE
    for j,(cell,cx,cw) in enumerate(zip(row,cxs2,cws2)):
        rect(s4,cx,ry,cw-Inches(0.04),Inches(0.58),fill=rbg,
             line_col=C_GRAY,line_w=Pt(0.5))
        tc2=GREEN if(j==3 and cell.startswith("+")) else(RED if(j==3 and cell.startswith("-")) else DARK)
        tb(s4,cell,cx+Inches(0.04),ry+Inches(0.06),cw,Inches(0.45),
           sz=10,color=tc2,align=PP_ALIGN.CENTER)

# Charts
roc=os.path.join(PI,"fig_roc_pr_curves.png")
if os.path.exists(roc):
    img(s4,roc,Inches(5.05),Inches(2.35),Inches(4.0),Inches(2.15))
thr=os.path.join(PI,"fig_threshold_sensitivity.png")
if os.path.exists(thr):
    img(s4,thr,Inches(9.2),Inches(2.35),Inches(3.85),Inches(2.15))
cm=os.path.join(PI,"fig_confusion_matrices.png")
if os.path.exists(cm):
    img(s4,cm,Inches(5.05),Inches(4.6),Inches(8.0),Inches(2.2))

# ════════════════════════════════════════════════════════════
# SLIDE 5 – CONCLUSION
# ════════════════════════════════════════════════════════════
s5=slide()
rect(s5,0,0,W,H,fill=BG)
hdr(s5,"Conclusion & Future Scope")

summ=[
    (C_GREEN,GREEN, "Achieved",
     "94.85% test accuracy  |  ROC-AUC 0.9899\n26,534 labeled samples across 8 diverse datasets  |  N_test = 3,980"),
    (C_BLUE, BLUE,  "Unique Value",
     "4-tier uncertainty routing + HITL loop absent in all surveyed systems\nHNM improved fraud recall +0.21pp with zero manual labeling cost"),
    (C_PURP, PURPLE,"Actionable Output",
     "First system to auto-generate NCRP-compliant PDFs + Velocity\nIntelligence DB giving law-enforcement threat tracking"),
    (C_AMBER,AMBER, "Efficiency",
     "End-to-end inference < 500ms on commodity hardware\nAll 5 functional + 5 non-functional requirements verified"),
]
for i,(bg2,col2,label,desc) in enumerate(summ):
    y=Inches(1.05)+i*Inches(1.22)
    card(s5,Inches(0.3),y,Inches(8.0),Inches(1.1),bg2,col2)
    dot3=s5.shapes.add_shape(9,Inches(0.48),y+Inches(0.26),Inches(0.2),Inches(0.2))
    dot3.fill.solid();dot3.fill.fore_color.rgb=col2;dot3.line.fill.background()
    tb(s5,label,Inches(0.78),y+Inches(0.06),Inches(2.0),Inches(0.38),
       sz=14,bold=True,color=col2)
    tb(s5,desc,Inches(0.78),y+Inches(0.46),Inches(7.3),Inches(0.58),
       sz=11,color=GRAY)

# Future scope
card(s5,Inches(8.55),Inches(1.0),Inches(4.5),Inches(5.5),C_PURP,PURPLE)
tb(s5,"FUTURE SCOPE",Inches(8.7),Inches(1.08),Inches(4.2),Inches(0.4),
   sz=13,bold=True,color=PURPLE)
future=[
    (GREEN, "Near","Direct NCRP REST API integration"),
    (GREEN, "Near","OTP/QR Code artifact extraction"),
    (AMBER, "Mid", "Multilingual NLP — Devanagari scripts"),
    (AMBER, "Mid", "Transformer (BERT) semantic embeddings"),
    (RED,   "Long","Federated learning across state cyber cells"),
    (RED,   "Long","Real-time threat intelligence feed API"),
]
for i,(col2,tag,desc) in enumerate(future):
    y=Inches(1.58)+i*Inches(0.82)
    rect(s5,Inches(8.7),y,Inches(0.85),Inches(0.28),fill=col2)
    tb(s5,tag,Inches(8.7),y-Inches(0.02),Inches(0.88),Inches(0.3),
       sz=9,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    tb(s5,desc,Inches(9.62),y,Inches(3.25),Inches(0.7),sz=11,color=DARK)

# ════════════════════════════════════════════════════════════
# SLIDE 6 – REFERENCES
# ════════════════════════════════════════════════════════════
s6=slide()
rect(s6,0,0,W,H,fill=BG)
hdr(s6,"References")
refs=[
    "[1] Almeida, T.A., Hidalgo, J.M.G., & Yamakami, A. (2011). Contributions to the study of SMS spam filtering. ACM DocEng, 259-262.",
    "[2] Settles, B. (2012). Active Learning. Synthesis Lectures on AI and Machine Learning. Morgan & Claypool.",
    "[3] Holzinger, A. (2016). Interactive machine learning for health informatics. Brain Informatics, 3(2), 119-131.",
    "[4] Chawla, N.V. et al. (2002). SMOTE: Synthetic Minority Over-sampling Technique. JAIR, 16, 321-357.",
    "[5] Roy, S. et al. (2020). A Deep Learning based AI Model for SMS Spam Detection. IJEAT, 9(3).",
    "[6] Pedregosa, F. et al. (2011). Scikit-learn: Machine Learning in Python. JMLR, 12, 2825-2830.",
    "[7] Purkait, S. (2014). Phishing counter measures and their effectiveness. Info. Mgmt. & Comp. Security, 22(5).",
    "[8] Ministry of Home Affairs, Govt. of India (2023). Annual Report on Cybercrime - NCRP Statistics 2023.",
    "[9] Friedl, J.E.F. (2006). Mastering Regular Expressions, 3rd ed. O'Reilly Media.",
    "[10] National Crime Records Bureau (2023). Crime in India Report - Cyber Crimes Chapter.",
]
for i,ref in enumerate(refs):
    card(s6,Inches(0.3),Inches(1.02)+i*Inches(0.6),
         Inches(12.6),Inches(0.54),
         C_BLUE if i%2==0 else WHITE,
         C_GRAY)
    tb(s6,ref,Inches(0.45),Inches(1.06)+i*Inches(0.6),
       Inches(12.3),Inches(0.48),sz=11,color=DARK)

out = r"d:\miniProject\Spam_reporting_portal\CyberShield_Presentation.pptx"
prs.save(out)
print("Saved:", out)
